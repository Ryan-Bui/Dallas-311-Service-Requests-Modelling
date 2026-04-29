import os
import sys
import json
import uuid
import re
from pathlib import Path
from neo4j import GraphDatabase
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv

# Setup paths
ROOT = Path(__file__).resolve().parents[1]
sys.path.append(str(ROOT))
load_dotenv(ROOT / ".env", override=True)

from inference.llm_factory import get_llm

URI = os.getenv("NEO4J_URI")
USERNAME = os.getenv("NEO4J_USERNAME")
PASSWORD = os.getenv("NEO4J_PASSWORD")

KNOWLEDGE_DIR = ROOT / "knowledge"

def sanitize_relationship_type(rel_text):
    if not isinstance(rel_text, str):
        return "RELATED_TO"
    clean = re.sub(r"[^a-zA-Z0-9\s_-]", "", rel_text)
    clean = re.sub(r"[\s-]+", "_", clean)
    clean = clean.upper().strip("_")
    return clean if clean else "RELATED_TO"


def parse_first_json_payload(text):
    """Return the first valid JSON object/array from an LLM response."""
    decoder = json.JSONDecoder()
    stripped = text.strip()

    # Remove common markdown fences without assuming the whole response is clean.
    stripped = re.sub(r"^```(?:json)?\s*", "", stripped, flags=re.IGNORECASE)
    stripped = re.sub(r"\s*```$", "", stripped)

    for index, char in enumerate(stripped):
        if char not in "[{":
            continue
        try:
            payload, _ = decoder.raw_decode(stripped[index:])
            return payload
        except json.JSONDecodeError:
            continue

    raise ValueError("No valid JSON object or array found in LLM response")


def normalize_triples(payload):
    """Coerce common LLM JSON shapes into a list of triple dictionaries."""
    if isinstance(payload, dict):
        for key in ("triples", "relationships", "data", "results"):
            value = payload.get(key)
            if isinstance(value, list):
                payload = value
                break
        else:
            payload = [payload]

    if not isinstance(payload, list):
        return []

    triples = []
    for item in payload:
        if isinstance(item, dict):
            triples.append(item)
            continue

        # Some models return a list of JSON strings. Accept those if they parse
        # cleanly to an object; otherwise skip them.
        if isinstance(item, str):
            try:
                parsed = parse_first_json_payload(item)
            except ValueError:
                continue
            if isinstance(parsed, dict):
                triples.append(parsed)
            elif isinstance(parsed, list):
                triples.extend(t for t in parsed if isinstance(t, dict))

    return triples

def extract_triples():
    if not all([URI, USERNAME, PASSWORD]):
        print("Error: Required Neo4j credentials missing in .env")
        return

    driver = GraphDatabase.driver(URI, auth=(USERNAME, PASSWORD))
    
    # Use the unified resilience factory LLM
    llm = get_llm(temperature=0.1)
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=350, chunk_overlap=50)

    # 1. Gather all files
    files = []
    if KNOWLEDGE_DIR.exists():
        for file in KNOWLEDGE_DIR.iterdir():
            if file.suffix in [".md", ".pptx"]:
                files.append(file)
                
    print(f"Found {len(files)} files to extract knowledge triples from.")

    for file_path in files:
        print(f"\nProcessing {file_path.name}...")
        extracted_text = ""
        
        if file_path.suffix == ".md":
            extracted_text = file_path.read_text(encoding="utf-8")
        elif file_path.suffix == ".pptx":
            try:
                prs = Presentation(file_path)
                slides_text = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        slides_text.append("\n".join(slide_text))
                extracted_text = "\n\n".join(slides_text)
            except Exception as e:
                print(f"  Error reading PPTX {file_path.name}: {e}")
                continue

        if not extracted_text.strip():
            print(f"  No text found in {file_path.name}.")
            continue

        chunks = text_splitter.split_text(extracted_text)
        print(f"  Created {len(chunks)} chunks.")

        for i, chunk in enumerate(chunks, 1):
            print(f"    Extracting triples from chunk {i}/{len(chunks)}...")
            prompt_template = """
You are extracting knowledge graph triples from Dallas City 311 service request data.

ENTITY TYPES you will encounter:
- Department         (e.g. "Dallas Water Utilities", "Code Compliance", "Street Services")
- ServiceRequestType (e.g. "Pothole Repair", "Graffiti Removal", "Missed Trash Pickup")
- CityCouncilDistrict (e.g. "District 1", "District 7")
- Priority           (e.g. "High", "Medium", "Low")
- Status             (e.g. "Open", "Closed", "In Progress")
- Outcome            (e.g. "Resolved", "No Action Required", "Duplicate")
- MethodReceived     (e.g. "Mobile App", "Phone", "Web Portal", "Walk-In")

HARD RULES:
- Never create ServiceRequest or Location entities.
- Never use one row, one request, a service request number, unique key, street address, or raw location as a node.
- Treat rows and locations only as evidence for relationships between ServiceRequestType, Department, CityCouncilDistrict, Priority, Outcome, Status, and MethodReceived.

RELATIONSHIP TAXONOMY (use ONLY these):
  HANDLED_BY         → ServiceRequestType  -[HANDLED_BY]->  Department
  LOCATED_IN         → ServiceRequest      -[LOCATED_IN]->  CityCouncilDistrict
  SUBMITTED_VIA      → ServiceRequest      -[SUBMITTED_VIA]-> MethodReceived
  ASSIGNED_PRIORITY  → ServiceRequest      -[ASSIGNED_PRIORITY]-> Priority
  RESOLVED_WITH      → ServiceRequest      -[RESOLVED_WITH]-> Outcome
  EXCEEDS_ERT        → ServiceRequest      -[EXCEEDS_ERT]->  Department  (when Closed Date > Due Date)
  MET_ERT            → ServiceRequest      -[MET_ERT]->      Department  (when Closed Date <= Due Date)
  RECURS_IN          → ServiceRequestType  -[RECURS_IN]->    CityCouncilDistrict  (same type appears 2+ times)
  COLOCATED_WITH     → ServiceRequestType  -[COLOCATED_WITH]-> ServiceRequestType  (same address, different type)
  ESCALATED_TO       → ServiceRequest      -[ESCALATED_TO]-> Department  (when Priority = High and status changed)
  OVERLAPS_WITH      → CityCouncilDistrict -[OVERLAPS_WITH]-> Department  (district consistently uses same dept)

CRITICAL RULE — EXTRACT ALL ENTITY PAIRS, NOT JUST ENTITY→HUB:

Given this 311 record:
  Department:          "Street Services"
  ServiceRequestType:  "Pothole Repair"
  CityCouncilDistrict: "District 4"
  Priority:            "High"
  Outcome:             "Resolved"
  Closed Date > Due Date: True

Extract ALL of these — not just each entity pointing at the request:
  ("Pothole Repair")   -[HANDLED_BY]->       ("Street Services")       ✅ lateral
  ("Pothole Repair")   -[RECURS_IN]->        ("District 4")            ✅ lateral
  ("Street Services")  -[OVERLAPS_WITH]->    ("District 4")            ✅ lateral
  ("Street Services")  -[EXCEEDS_ERT]->      ("High Priority")         ✅ performance signal
  ("Pothole Repair")   -[ASSIGNED_PRIORITY]-> ("High")                 ✅ 
  ("District 4")       -[RESOLVED_WITH]->    ("Resolved")              ✅

NEVER extract:
  (anything) -[BELONGS_TO]->   (ServiceRequest_12345)   ❌ request as hub
  (anything) -[MENTIONED_IN]-> (Record)                 ❌ document as hub
  (anything) -[RELATED_TO]->   (anything)               ❌ too vague
  (anything) -[KNOWLEDGE_REL]-> (anything)              ❌ banned

LATERAL EXTRACTION RULE:
For every batch of records, if you find entities A, B, and C:
  → Extract A→B, A→C, AND B→C
  → Never only extract A→hub, B→hub, C→hub
  → The most valuable edges are between two non-request entities

PERFORMANCE PATTERN EXTRACTION (high value for RAG):
Look across multiple records and extract:
  - Which departments consistently exceed ERT for which request types
  - Which districts have recurring request types (RECURS_IN)
  - Which methods of submission correlate with faster closure
  - Which outcomes are associated with high priority requests

Example multi-record inference:
  Records show: District 7 has 14 "Illegal Dumping" requests, all handled by 
  "Code Compliance", 9 of 14 exceeded ERT.
  
  Extract:
  ("Illegal Dumping")   -[RECURS_IN]->      ("District 7")         ✅
  ("Code Compliance")   -[EXCEEDS_ERT_FOR]-> ("Illegal Dumping")   ✅
  ("District 7")        -[OVERLOADS]->       ("Code Compliance")   ✅

OUTPUT FORMAT:
Return ONLY a JSON array. No preamble, no markdown.
[
  {
    "subject":      "entity name",
    "subject_type": "EntityType",
    "predicate":    "RELATIONSHIP_TYPE",
    "object":       "entity name",
    "object_type":  "EntityType",
    "confidence":   0.0-1.0,
    "evidence":     "the specific field values that support this triple"
  }
]

Drop any triple where confidence < 0.75.
Drop any triple where evidence is empty.
Drop any triple where subject or object is a Service Request Number, Unique Key, street address, raw location, or one individual record/request.

Text: "{chunk}"
"""
            prompt = prompt_template.replace("{chunk}", chunk)
            
            try:
                response = llm.invoke(prompt)
                res_text = response.content.strip()
                
                res_json = parse_first_json_payload(res_text)
                triples = normalize_triples(res_json)

                if not triples:
                    print(f"      No triples extracted.")
                    continue

                print(f"      Ingesting {len(triples)} triples...")
                with driver.session() as session:
                    for t in triples:
                        sub = str(t.get("subject", "")).strip()
                        rel = str(t.get("predicate", t.get("relation", ""))).strip()
                        obj = str(t.get("object", "")).strip()
                        sub_type = str(t.get("subject_type", "")).strip().lower()
                        obj_type = str(t.get("object_type", "")).strip().lower()
                        
                        if not sub or not obj or not rel:
                            continue

                        blocked_types = {"servicerequest", "service request", "location"}
                        if sub_type in blocked_types or obj_type in blocked_types:
                            continue

                        blocked_value = re.compile(
                            r"(^|\b)(sr[_\s-]*\d+|service request|unique key|record\s+\d+|request\s+\d+)(\b|$)",
                            re.IGNORECASE,
                        )
                        if blocked_value.search(sub) or blocked_value.search(obj):
                            continue
                            
                        sub_id = str(uuid.uuid5(uuid.NAMESPACE_DNS, sub.lower()))
                        obj_id = str(uuid.uuid5(uuid.NAMESPACE_DNS, obj.lower()))
                        
                        rel_type = sanitize_relationship_type(rel)
                        session.run(f"""
                            MERGE (s:KnowledgeEntity {{id: $sub_id}}) ON CREATE SET s.name = $sub_name
                            MERGE (o:KnowledgeEntity {{id: $obj_id}}) ON CREATE SET o.name = $obj_name
                            MERGE (s)-[r:{rel_type}]->(o)
                            SET r.relation = $rel, r.source = $source
                        """, sub_id=sub_id, sub_name=sub, obj_id=obj_id, obj_name=obj, rel=rel, source=file_path.name)
                        
            except Exception as e:
                print(f"      Error extracting/ingesting: {e}")

    driver.close()
    print("\n[SUCCESS] Entity-Relationship triples stored securely in Neo4j Knowledge Graph!")

if __name__ == "__main__":
    extract_triples()
